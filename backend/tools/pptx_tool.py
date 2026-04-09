import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from config import settings


def generate_pptx(title: str, slides_content: list[dict], brand: dict, filename: str) -> str:
    """
    slides_content: list of {"title": str, "body": str}
    """
    os.makedirs(settings.output_dir, exist_ok=True)
    prs = Presentation()
    primary_color = _hex_to_rgb(brand.get("primary_color", "#000000"))
    font = brand.get("font", "Calibri")

    for slide_data in slides_content:
        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)

        # Title
        tf = slide.shapes.title.text_frame
        tf.text = slide_data.get("title", "")
        tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(*primary_color)
        tf.paragraphs[0].runs[0].font.name = font
        tf.paragraphs[0].runs[0].font.size = Pt(28)

        # Body
        body = slide.placeholders[1].text_frame
        body.text = slide_data.get("body", "")
        for para in body.paragraphs:
            for run in para.runs:
                run.font.name = font
                run.font.size = Pt(16)

    path = os.path.join(settings.output_dir, filename)
    prs.save(path)
    return path


def _hex_to_rgb(hex_color: str) -> tuple:
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
